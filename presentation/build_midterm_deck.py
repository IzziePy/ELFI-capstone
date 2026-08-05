"""Build the ELFI midterm presentation (English) as a .pptx file."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ------------------------------------------------------------------
# Design tokens
# ------------------------------------------------------------------

INK = RGBColor(0x0F, 0x1A, 0x1C)
MUTED = RGBColor(0x5A, 0x6B, 0x6E)
FAINT = RGBColor(0x8A, 0x9B, 0x9E)
ACCENT = RGBColor(0x14, 0x63, 0x6E)
ACCENT_LT = RGBColor(0xE2, 0xF0, 0xF1)
SUNKEN = RGBColor(0xEE, 0xF2, 0xF2)
RULE = RGBColor(0xD0, 0xDA, 0xDB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
WARM = RGBColor(0x8F, 0x64, 0x14)
CRIT = RGBColor(0x9A, 0x40, 0x37)
OK = RGBColor(0x2A, 0x66, 0x4C)

SERIF = "Georgia"
SANS = "Calibri"
MONO = "Consolas"

SW, SH = 13.333, 7.5
M = 0.82                      # side margin
CW = SW - 2 * M               # content width
BODY_TOP = 2.05
BODY_BOTTOM = 6.72

prs = Presentation()
prs.slide_width = Inches(SW)
prs.slide_height = Inches(SH)
BLANK = prs.slide_layouts[6]


# ------------------------------------------------------------------
# Primitives
# ------------------------------------------------------------------

def textbox(slide, x, y, w, h):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def para(tf, text, *, size=12, font=SANS, color=INK, bold=False, italic=False,
         space_before=0, space_after=4, line=1.2, caps=False, align=None, first=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    r = p.add_run()
    r.text = text.upper() if caps else text
    f = r.font
    f.size = Pt(size)
    f.name = font
    f.color.rgb = color
    f.bold = bold
    f.italic = italic
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    p.line_spacing = line
    if align:
        p.alignment = align
    return p


def bullet(tf, text, *, size=11.5, color=INK, first=False, space_after=4,
           mark="▪", bold_prefix=None, line=1.18):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    m = p.add_run()
    m.text = mark + "  "
    m.font.size = Pt(size)
    m.font.name = SANS
    m.font.color.rgb = ACCENT
    if bold_prefix:
        b = p.add_run()
        b.text = bold_prefix
        b.font.size = Pt(size)
        b.font.name = SANS
        b.font.bold = True
        b.font.color.rgb = INK
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.name = SANS
    r.font.color.rgb = color
    p.space_after = Pt(space_after)
    p.line_spacing = line
    return p


def rect(slide, x, y, w, h, fill, line=None, shape=MSO_SHAPE.RECTANGLE, radius=None):
    sh = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        sh.adjustments[0] = radius
    sh.text_frame.word_wrap = True
    return sh


def card(slide, x, y, w, h, *, lift=False):
    """Panel with a left accent keyline."""
    fill = ACCENT_LT if lift else SUNKEN
    rect(slide, x, y, w, h, fill, RULE)
    rect(slide, x, y, 0.035, h, ACCENT)
    tf = textbox(slide, x + 0.28, y + 0.22, w - 0.52, h - 0.44)
    return tf


def label(tf, text, *, first=False, color=ACCENT, size=9.5, space_after=6):
    return para(tf, text, size=size, font=MONO, color=color, bold=True,
                caps=True, space_after=space_after, first=first)


def slide_frame(title, eyebrow, kicker=None, *, section="ELFI · Midterm"):
    slide = prs.slides.add_slide(BLANK)
    rect(slide, 0, 0, 0.055, SH, ACCENT)

    tf = textbox(slide, M, 0.46, CW, 0.3)
    p = para(tf, eyebrow, size=10, font=MONO, color=ACCENT, bold=True,
             caps=True, space_after=0, first=True)
    if kicker:
        r = p.add_run()
        r.text = "   ·   " + kicker.upper()
        r.font.size = Pt(10)
        r.font.name = MONO
        r.font.color.rgb = FAINT

    tf = textbox(slide, M, 0.86, CW * 0.86, 1.0)
    para(tf, title, size=30, font=SERIF, color=INK, space_after=0, line=1.0, first=True)

    tf = textbox(slide, M, 6.95, CW, 0.28)
    p = para(tf, section, size=9, font=MONO, color=FAINT, caps=True,
             space_after=0, first=True)
    slide._elfi_footer = p
    return slide


def stamp_page(slide, n, total):
    tf = textbox(slide, SW - M - 2.0, 6.95, 2.0, 0.28)
    para(tf, f"{n} / {total}", size=9, font=MONO, color=FAINT,
         space_after=0, align=PP_ALIGN.RIGHT, first=True)


def notes(slide, *paragraphs):
    tf = slide.notes_slide.notes_text_frame
    tf.text = paragraphs[0]
    for extra in paragraphs[1:]:
        tf.add_paragraph().text = extra


def table(slide, x, y, w, rows, *, col_ratios=None, header=True,
          row_h=0.29, header_h=0.3, size=10.5, right_align=()):
    n_rows, n_cols = len(rows), len(rows[0])
    total_h = header_h + row_h * (n_rows - 1) if header else row_h * n_rows
    shape = slide.shapes.add_table(n_rows, n_cols, Inches(x), Inches(y),
                                   Inches(w), Inches(total_h))
    tbl = shape.table
    tbl.first_row = header
    tbl.horz_banding = False

    ratios = col_ratios or [1 / n_cols] * n_cols
    for c, ratio in enumerate(ratios):
        tbl.columns[c].width = Inches(w * ratio)
    for r in range(n_rows):
        tbl.rows[r].height = Inches(header_h if (header and r == 0) else row_h)

    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.margin_left = Inches(0.07)
            cell.margin_right = Inches(0.07)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            if header and r == 0:
                cell.fill.fore_color.rgb = ACCENT_LT
            else:
                cell.fill.fore_color.rgb = WHITE if r % 2 else SUNKEN

            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            text, bold, colour, mono = val, False, INK, False
            if isinstance(val, tuple):
                text = val[0]
                opts = val[1] if len(val) > 1 else {}
                bold = opts.get("bold", False)
                colour = opts.get("color", INK)
                mono = opts.get("mono", False)
            r_ = p.add_run()
            r_.text = str(text)
            f = r_.font
            f.size = Pt(size if not (header and r == 0) else size - 1)
            f.name = MONO if (mono or (header and r == 0)) else SANS
            f.bold = bold or (header and r == 0)
            f.color.rgb = MUTED if (header and r == 0) else colour
            p.line_spacing = 1.0
            if c in right_align:
                p.alignment = PP_ALIGN.RIGHT
    return tbl


def kpi(slide, x, y, w, value, key, note=None, *, lift=False):
    h = 1.55 if note else 1.05
    fill = ACCENT_LT if lift else SUNKEN
    rect(slide, x, y, w, h, fill, RULE)
    tf = textbox(slide, x + 0.24, y + 0.2, w - 0.48, h - 0.34)
    para(tf, value, size=26, font=SERIF, color=ACCENT if lift else INK,
         space_after=2, line=1.0, first=True)
    para(tf, key, size=9, font=MONO, color=MUTED, caps=True, space_after=4)
    if note:
        para(tf, note, size=10, color=MUTED, line=1.15, space_after=0)


def pill(slide, x, y, w, text, color):
    sh = rect(slide, x, y, w, 0.24, None, color, MSO_SHAPE.ROUNDED_RECTANGLE, 0.3)
    tf = sh.text_frame
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = 0
    para(tf, text, size=8.5, font=MONO, color=color, bold=True, caps=True,
         space_after=0, align=PP_ALIGN.CENTER, first=True)
    return sh


# ------------------------------------------------------------------
# 1 — Title
# ------------------------------------------------------------------

s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, 0.11, ACCENT)

tf = textbox(s, M, 1.15, 8.5, 1.8)
para(tf, "ELFI", size=88, font=SERIF, color=INK, space_after=8, line=0.92, first=True)
para(tf, "Emotional & Latent Feedback Interpretation", size=17, font=MONO,
     color=ACCENT, caps=True, space_after=0, line=1.1)

rect(s, M, 3.62, CW, 0.01, RULE)

tf = textbox(s, M, 3.95, 9.6, 2.1)
for i, q in enumerate(["How do employees feel?",
                       "What are employees talking about?",
                       "How do emotions relate to engagement?"]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    n = p.add_run()
    n.text = f"0{i + 1}   "
    n.font.size = Pt(11)
    n.font.name = MONO
    n.font.color.rgb = ACCENT
    n.font.bold = True
    r = p.add_run()
    r.text = q
    r.font.size = Pt(20)
    r.font.name = SERIF
    r.font.color.rgb = MUTED
    p.space_after = Pt(9)
    p.line_spacing = 1.1

tf = textbox(s, M, 6.6, CW, 0.4)
para(tf, "Capstone midterm     ·     NLP · Core Affect · Power BI"
         "     ·     Data: High5 employee survey",
     size=10, font=MONO, color=FAINT, caps=True, space_after=0, first=True)

notes(s,
      "Good morning. My project is called ELFI — Emotional & Latent Feedback "
      "Interpretation.",
      "Companies collect employee feedback regularly, and a large part of it arrives as "
      "free-text comments. Those comments carry a lot of valuable information, but they "
      "are usually only partly used.",
      "There is a second problem. Smaller units typically read only their own comments and "
      "then draw conclusions about the whole organisation. The cross-organisational picture "
      "is missing.",
      "So my goal is a dashboard that automatically detects which topics are discussed, "
      "which emotions sit behind them, and how both shift between business units and over "
      "time.")

# ------------------------------------------------------------------
# 2 — Problem
# ------------------------------------------------------------------

s = slide_frame("Thousands of comments are collected. Few of them are used.",
                "Problem", "Why this project exists")

col_w = (CW - 2 * 0.32) / 3
blocks = [
    ("Read, not analysed",
     "Free-text answers are skimmed manually, quoted selectively and then archived. "
     "There is no systematic, repeatable read of the corpus."),
    ("Local view, global claim",
     "Each unit sees only its own comments and generalises from them. Whether an issue "
     "is local or organisation-wide cannot be answered."),
    ("Scores without a why",
     "Likert scores show that a dimension dropped. The comments hold the explanation — "
     "but they are never linked to the scores."),
]
for i, (head, body) in enumerate(blocks):
    tf = card(s, M + i * (col_w + 0.32), BODY_TOP, col_w, 2.0)
    label(tf, head, first=True)
    para(tf, body, size=11.5, color=MUTED, line=1.25, space_after=0)

tf = card(s, M, 4.4, CW, 1.5, lift=True)
label(tf, "Guiding question", first=True)
para(tf, "How can unused employee feedback become a strategic management instrument?",
     size=21, font=SERIF, color=INK, line=1.15, space_after=0)

notes(s,
      "Three things are wrong with the status quo: comments are read but never analysed "
      "systematically; every unit generalises from its own small sample; and the survey "
      "scores tell us that something moved without telling us why.",
      "ELFI attacks all three by turning the comment corpus into a structured, comparable "
      "data layer that can be sliced by unit and by iteration.")

# ------------------------------------------------------------------
# 3 — Why emotions
# ------------------------------------------------------------------

s = slide_frame("Negative is not one thing.", "Framework", "Why emotions, not polarity")

half = (CW - 0.4) / 2
tf = card(s, M, BODY_TOP, half, 4.2)
label(tf, "Classic sentiment", first=True)
for cls in ["positive", "neutral", "negative"]:
    para(tf, cls, size=12, font=MONO, color=MUTED, caps=True, space_after=7)
para(tf, "One negative comment may express anger. Another expresses uncertainty, "
         "a third exhaustion.", size=12, color=INK, line=1.3, space_before=6, space_after=6)
para(tf, "Same label — completely different measures.", size=12, color=MUTED,
     italic=True, line=1.3, space_after=10)
bullet(tf, "Three classes cannot separate an angry team from a team that has quietly "
           "given up", size=11)
bullet(tf, "One share of negatives per unit is a number, not a direction for action",
       size=11)
bullet(tf, "Intensity is lost: mild irritation and open conflict share a label",
       size=11, space_after=0)

x2 = M + half + 0.4
tf = card(s, x2, BODY_TOP, half, 4.2)
label(tf, "Core Affect — valence × arousal", first=True)

# quadrant diagram
qx, qy, qw, qh = x2 + 0.34, BODY_TOP + 0.78, half - 0.68, 2.35
rect(s, qx, qy, qw / 2, qh / 2, RGBColor(0xF7, 0xFA, 0xFA), None)
rect(s, qx + qw / 2, qy, qw / 2, qh / 2, RGBColor(0xE9, 0xF3, 0xF4), None)
rect(s, qx, qy + qh / 2, qw / 2, qh / 2, RGBColor(0xF7, 0xFA, 0xFA), None)
rect(s, qx + qw / 2, qy + qh / 2, qw / 2, qh / 2, RGBColor(0xF2, 0xF7, 0xF7), None)
rect(s, qx, qy + qh / 2 - 0.006, qw, 0.012, RULE)
rect(s, qx + qw / 2 - 0.006, qy, 0.012, qh, RULE)

quad = [("Stress", "tension, anger, overload", 0, 0),
        ("Enthusiasm", "energy, pride, momentum", 1, 0),
        ("Resignation", "withdrawal, fatigue", 0, 1),
        ("Satisfaction", "calm, trust, stability", 1, 1)]
for name, sub, cx, cy in quad:
    tf_q = textbox(s, qx + cx * qw / 2 + 0.12, qy + cy * qh / 2 + 0.12, qw / 2 - 0.24, 0.6)
    para(tf_q, name, size=11.5, color=INK, bold=True, space_after=1, first=True)
    para(tf_q, sub, size=9, color=MUTED, space_after=0, line=1.1)

tf_a = textbox(s, qx, qy + qh + 0.04, qw / 2, 0.24)
para(tf_a, "← valence −", size=8.5, font=MONO, color=FAINT, caps=True,
     space_after=0, first=True)
tf_a = textbox(s, qx + qw / 2, qy + qh + 0.04, qw / 2, 0.24)
para(tf_a, "valence + →", size=8.5, font=MONO, color=FAINT, caps=True,
     space_after=0, align=PP_ALIGN.RIGHT, first=True)
tf_a = textbox(s, qx, qy - 0.26, qw, 0.24)
para(tf_a, "↑ arousal + (activated)", size=8.5, font=MONO, color=FAINT, caps=True,
     space_after=0, first=True)

tf = textbox(s, x2 + 0.34, 5.62, half - 0.68, 0.6)
para(tf, "After Lisa Feldman Barrett's Theory of Constructed Emotion: emotions are not "
         "fixed categories but constructed over two continuous dimensions.",
     size=11, color=MUTED, line=1.25, space_after=0, first=True)

notes(s,
      "Most sentiment systems work with positive, neutral and negative. For employee "
      "feedback that is too coarse. A negative comment can point to anger, but it can just "
      "as well express uncertainty or feeling overwhelmed — and those require completely "
      "different measures.",
      "That is why I follow Lisa Feldman Barrett's Theory of Constructed Emotion: emotions "
      "are described as a combination of valence and arousal instead of fixed classes.",
      "Example: uncertainty during a change process argues for additional communication and "
      "workshops. Anger about workload argues for something else entirely.")

# ------------------------------------------------------------------
# 4 — Data & privacy
# ------------------------------------------------------------------

s = slide_frame("Developed on real data. Demonstrated on synthetic data.",
                "Data", "Two datasets, one pipeline")

tf = card(s, M, BODY_TOP, half, 3.05)
label(tf, "Real High5 data — confidential", first=True)
table(s, M + 0.28, BODY_TOP + 0.62, half - 0.56,
      [["Responses", ("46,000+", {"bold": True})],
       ["Free-text comments", ("4,297", {"bold": True})],
       ["Likert items", ("54", {"bold": True})],
       ["Survey exports", ("17 (23_1 → 24_2)", {"bold": True})]],
      header=False, col_ratios=[0.58, 0.42], right_align=(1,), row_h=0.3)
tf2 = textbox(s, M + 0.28, BODY_TOP + 2.05, half - 0.56, 1.6)
bullet(tf2, "8 thematic High5 dimensions, several business units", first=True)
bullet(tf2, "Statistical figures may be shown — original text may not")
bullet(tf2, "All model development happens here", space_after=0)

tf = card(s, x2, BODY_TOP, half, 3.05)
label(tf, "Synthetic data — publishable", first=True)
table(s, x2 + 0.28, BODY_TOP + 0.62, half - 0.56,
      [["Responses / comments", ("7,000 / 5,000", {"bold": True})],
       ["Comment rate", ("> 96 %", {"bold": True})],
       ["Median comment length", ("25 words", {"bold": True})],
       ["Cronbach α per dimension", ("0.68 – 0.96", {"bold": True})]],
      header=False, col_ratios=[0.58, 0.42], right_align=(1,), row_h=0.3)
tf2 = textbox(s, x2 + 0.28, BODY_TOP + 2.05, half - 0.56, 1.6)
bullet(tf2, "Same schema, same scales, same iteration structure", first=True)
bullet(tf2, "Dimension scores 3.1–3.7 of 5; r(teamwork, culture) ≈ 0.89")
bullet(tf2, "Commenters rate more critically — as in the real data", space_after=0)

tf = card(s, M, 5.4, CW, 1.3)
label(tf, "The privacy rule", first=True)
para(tf, "Every model runs on both datasets. Development and evaluation happen on the "
         "confidential data; everything that leaves the company — screenshots, notebooks, "
         "the published dashboard — uses the synthetic twin. Statistical figures from the "
         "real data may be quoted, comment text may not.",
     size=12, color=INK, line=1.3, space_after=0)

notes(s,
      "Development happens on the real High5 data. The comments are confidential, so I "
      "cannot show original text — but I am allowed to use statistical figures.",
      "For the presentation and for later publication I therefore use a synthetic dataset: "
      "same structures as the original, similar statistical properties, no confidential "
      "content.",
      "In the EDA I examined data quality, missing values, distributions, iterations and the "
      "scoring structure, as well as the differences between the real and the synthetic "
      "dataset.")

# ------------------------------------------------------------------
# 5 — Architecture
# ------------------------------------------------------------------

s = slide_frame("The pipeline", "Architecture", "Survey data to dashboard")

steps = [("01", "Survey data", "17 exports, harmonised schema"),
         ("02", "EDA", "Quality, missingness, reliability"),
         ("03", "Cleaning", "Recoding, reverse coding, scores"),
         ("04", "Sentiment", "German Sentiment BERT → valence"),
         ("05", "Topics", "BERTopic → clusters, trends"),
         ("06", "Dashboard", "Power BI, enriched export")]
sw_ = (CW - 5 * 0.16) / 6
for i, (n, t, d) in enumerate(steps):
    x = M + i * (sw_ + 0.16)
    rect(s, x, BODY_TOP, sw_, 1.24, SUNKEN, RULE)
    rect(s, x, BODY_TOP, 0.035, 1.24, ACCENT)
    tf = textbox(s, x + 0.18, BODY_TOP + 0.15, sw_ - 0.34, 1.0)
    para(tf, n, size=9, font=MONO, color=FAINT, space_after=3, first=True)
    para(tf, t, size=12.5, color=INK, bold=True, space_after=3, line=1.05)
    para(tf, d, size=9.5, color=MUTED, space_after=0, line=1.15)
    if i < 5:
        tf_a = textbox(s, x + sw_ - 0.02, BODY_TOP + 0.47, 0.2, 0.3)
        para(tf_a, "›", size=14, font=MONO, color=FAINT, space_after=0,
             align=PP_ALIGN.CENTER, first=True)

tf = card(s, M, 3.6, half, 2.25)
label(tf, "Design decisions", first=True)
bullet(tf, "Emotion and topic are modelled separately, then joined per comment — each "
           "layer stays interpretable on its own.", bold_prefix="Two independent layers.  ")
bullet(tf, "The comments are German; multilingual defaults lose nuance.",
       bold_prefix="German-specific models.  ")
bullet(tf, "The pipeline ends in a flat, dashboard-ready dataset, not in a notebook.",
       bold_prefix="One analytical table.  ", space_after=0)

tf = card(s, x2, 3.6, half, 2.25)
label(tf, "In evaluation — caching", first=True)
para(tf, "Embeddings, sentiment inference and BERTopic fits are the expensive steps. "
         "I am currently evaluating persisting intermediate results as Parquet files so "
         "they are not recomputed on every run.",
     size=11.5, color=INK, line=1.3, space_after=6)
para(tf, "Evaluated, not yet implemented — details in the appendix.",
     size=11.5, color=MUTED, italic=True, space_after=0)

tf = textbox(s, M, 6.05, CW, 0.4)
para(tf, "One full pass over 5,000 comments takes about five minutes for the sentiment "
         "and topic layers together — the cost sits in tuning, not in production.",
     size=11.5, color=MUTED, italic=True, space_after=0, line=1.25, first=True)

notes(s,
      "The technical pipeline has several stages. First exploratory data analysis and "
      "cleaning. Then the comments are analysed semantically: German Sentiment BERT for the "
      "emotional layer, BERTopic for the thematic layer. Both results are prepared for the "
      "dashboard.",
      "To speed up the pipeline I am also evaluating storing intermediate results as Parquet "
      "files, so the compute-heavy steps do not have to run again every time.")

# ------------------------------------------------------------------
# 6 — Models & status
# ------------------------------------------------------------------

s = slide_frame("Three layers, one dataset", "Models & status", "What is built")

tf = card(s, M, BODY_TOP, half, 4.2)
label(tf, "Model stack", first=True)
table(s, M + 0.28, BODY_TOP + 0.62, half - 0.56,
      [["Layer", "Model", "Output"],
       ["Sentiment", ("german-sentiment-bert", {"mono": True}), "continuous valence"],
       ["Topics", ("BERTopic + MiniLM", {"mono": True}), "topic id, label, prob."],
       ["Baseline", ("mDeBERTa zero-shot", {"mono": True}), "valence / arousal class"],
       ["Dashboard", ("Power BI", {"mono": True}), "4 report pages"]],
      col_ratios=[0.22, 0.44, 0.34], size=10)
tf2 = textbox(s, M + 0.28, BODY_TOP + 2.35, half - 0.56, 1.4)
para(tf2, "Valence = P(positive) − P(negative), continuous on −1 … +1.",
     size=11.5, font=MONO, color=ACCENT, space_after=6, first=True)
para(tf2, "BERTopic uses sentence embeddings, so semantically similar comments cluster "
          "even when they share no keywords.", size=11.5, color=MUTED, line=1.25,
     space_after=0)

tf = card(s, x2, BODY_TOP, half, 4.2)
label(tf, "Progress", first=True)
rows = [("EDA — real & synthetic", "done"),
        ("Cleaning & survey scores", "done"),
        ("Zero-shot proof of concept", "done"),
        ("Sentiment — continuous valence", "done"),
        ("BERTopic incl. grid searches", "done"),
        ("Dashboard concept", "done"),
        ("Arousal dimension", "open"),
        ("Engagement prediction model", "open")]
table(s, x2 + 0.28, BODY_TOP + 0.62, half - 0.56,
      [[name, (state.upper(), {"mono": True, "bold": True,
                               "color": OK if state == "done" else WARM})]
       for name, state in rows],
      header=False, col_ratios=[0.72, 0.28], right_align=(1,), row_h=0.29, size=10.5)
tf2 = textbox(s, x2 + 0.28, BODY_TOP + 3.05, half - 0.56, 0.9)
para(tf2, "The core building blocks of the MVP are implemented. What remains is the second "
          "emotional axis and the predictive layer.", size=11.5, color=INK, line=1.3,
     space_after=0, first=True)

notes(s,
      "For sentiment analysis I use German Sentiment BERT and turn its class probabilities "
      "into a continuous valence value.",
      "For the thematic analysis I use BERTopic, which relies on modern embeddings and works "
      "well for semantic analysis of German free text.",
      "Current state: EDA, cleaning, sentiment analysis, BERTopic and the dashboard concept "
      "are done. The two open pieces are the arousal axis and the predictive model.")

# ------------------------------------------------------------------
# 7 — Metrics
# ------------------------------------------------------------------

s = slide_frame("Topic Coherence (C_v)", "Metrics", "How success is measured")

tf = textbox(s, M, 1.92, CW * 0.7, 0.5)
para(tf, "Measures how well the terms inside one topic belong together semantically.",
     size=16, font=SERIF, color=MUTED, space_after=0, line=1.2, first=True)

kw = (CW - 3 * 0.28) / 4
kpi(s, M, 2.6, kw, "0.350", "C_v coherence",
    "Current best configuration. Higher = more consistent, more interpretable topics.",
    lift=True)
kpi(s, M + (kw + 0.28), 2.6, kw, "18.1 %", "outlier rate",
    "Share of comments in the noise cluster (−1). Coverage metric.")
kpi(s, M + 2 * (kw + 0.28), 2.6, kw, "100 %", "comment coverage",
    "All 5,000 valid comments received a valence score. Target ≥ 95 %.")
kpi(s, M + 3 * (kw + 0.28), 2.6, kw, "MAE · R²", "predictive layer",
    "Planned for engagement prediction. Targets MAE < 0.60, R² > 0.30.")

tf = card(s, M, 4.45, half, 1.15)
label(tf, "Coherent topic", first=True, space_after=4)
para(tf, "communication · feedback · transparency · information",
     size=12, font=MONO, color=OK, space_after=0)
tf = card(s, x2, 4.45, half, 1.15)
label(tf, "Incoherent topic", first=True, space_after=4)
para(tf, "feedback · summer · dog · finance",
     size=12, font=MONO, color=CRIT, space_after=0)

tf = textbox(s, M, 5.85, CW, 0.7)
para(tf, "Topic modelling has no ground-truth labels, so classification metrics do not "
         "apply. Coherence, outlier rate and coverage are reported together — each one "
         "alone can be gamed.", size=11.5, color=MUTED, line=1.25, space_after=0, first=True)

notes(s,
      "Topic modelling needs a different evaluation logic than classic classification: "
      "there is no ground truth.",
      "I currently use C_v coherence. It measures how well the terms within a topic fit "
      "together. A good topic contains semantically related terms; a bad topic is a mix of "
      "unrelated terms. The higher the value, the more consistent and interpretable the "
      "topics.",
      "Alongside coherence I track the outlier rate as a coverage metric, and comment "
      "coverage for the pipeline itself. The predictive layer will add MAE, RMSE and R².")

# ------------------------------------------------------------------
# 8 — Dashboard & roadmap
# ------------------------------------------------------------------

s = slide_frame("Four pages, one question each", "Product", "Target dashboard & roadmap")

pages = [("Page 1", "Overview",
          "Comment volume, average valence, participation, iteration filter.", False),
         ("Page 2", "Core Affect Map",
          "Comments and topic clusters positioned on valence × arousal. The centrepiece.",
          True),
         ("Page 3", "Topics & Trends",
          "Topic frequency, movement across iterations, emotion per topic.", False),
         ("Page 4", "Organisational Comparison",
          "Business units side by side, valence against survey scores.", False)]
for i, (tag, title_, body, lift) in enumerate(pages):
    x = M + i * (kw + 0.28)
    tf = card(s, x, BODY_TOP, kw, 2.5, lift=lift)
    para(tf, tag, size=9, font=MONO, color=FAINT, caps=True, space_after=4, first=True)
    para(tf, title_, size=15, font=SERIF, color=ACCENT if lift else INK,
         space_after=6, line=1.05)
    para(tf, body, size=11, color=MUTED, line=1.25, space_after=0)

tf = card(s, M, 4.85, CW, 1.4, lift=True)
label(tf, "Project progress ≈ 65 %", first=True)
para(tf, "Data layer, sentiment layer, topic layer and dashboard concept are in place. "
         "Remaining: the arousal axis, topic consolidation, the predictive model, dashboard "
         "finalisation and the evaluation write-up.",
     size=12.5, color=INK, line=1.3, space_after=0)

notes(s,
      "The central component of the dashboard will be the Core Affect Map: comments or "
      "topic clusters positioned by valence and arousal.",
      "Around it sit topic trends, topic development across iterations, comparisons between "
      "business units, and the relationship between emotions and survey scores.",
      "Overall project progress is at roughly two thirds.")

# ------------------------------------------------------------------
# 9 — Risks & next steps
# ------------------------------------------------------------------

s = slide_frame("What could go wrong, and what happens next", "Outlook",
                "Risks & next steps")

tf = card(s, M, BODY_TOP, half, 2.5)
label(tf, "Risks", first=True, color=CRIT)
bullet(tf, "Generated comments are more repetitive than real ones, which depresses topic "
           "quality — the models look worse than they are.",
       bold_prefix="Synthetic text quality.  ")
bullet(tf, "The second Core Affect axis has no validated model yet.",
       bold_prefix="Arousal still open.  ")
bullet(tf, "303 micro-topics with a median of 11 comments each is too granular to act on.",
       bold_prefix="Topic granularity.  ")
bullet(tf, "Full reruns are slow without caching.", bold_prefix="Compute cost.  ",
       space_after=0)

tf = card(s, x2, BODY_TOP, half, 2.5)
label(tf, "Next steps", first=True)
bullet(tf, "— a German emotion model mapped onto valence and arousal coordinates",
       bold_prefix="Emotion mapping ")
bullet(tf, "— merge micro-topics into actionable themes",
       bold_prefix="Topic consolidation ")
bullet(tf, "across iterations, and organisational comparisons",
       bold_prefix="Topic trends  ")
bullet(tf, "in Power BI", bold_prefix="Dashboard finalisation  ")
bullet(tf, "— coherence, valence vs. engagement, predictive metrics",
       bold_prefix="Evaluation ", space_after=0)

rect(s, M, 4.95, 0.03, 1.15, ACCENT)
tf = textbox(s, M + 0.28, 5.0, CW - 0.5, 1.15)
para(tf, "My goal is to build a system that turns thousands of open employee comments into "
         "more than a list of topics: it shows which emotions sit behind them, and which "
         "measures can be derived from that — for a single unit or for the whole "
         "organisation.",
     size=15, font=SERIF, color=INK, line=1.25, space_after=0, first=True)

notes(s,
      "Name the risks openly — examiners reward that more than a clean story.",
      "Then close with the sentence on the slide, word for word, and stop. The appendix is "
      "there for the questions.")

# ------------------------------------------------------------------
# 10 — Appendix divider
# ------------------------------------------------------------------

s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, SUNKEN)
rect(s, 0, 0, SW, 0.11, ACCENT)
tf = textbox(s, M, 2.7, 9.5, 2.0)
para(tf, "Appendix", size=54, font=SERIF, color=INK, space_after=10, line=1.0, first=True)
para(tf, "Technical deep dive — backup slides, not presented",
     size=14, font=MONO, color=ACCENT, caps=True, space_after=12)
para(tf, "One slide per pipeline stage: what was analysed, which technical steps were "
         "taken, why, and what came out. Written to be read, not shown.",
     size=13, color=MUTED, line=1.35, space_after=0)
notes(s, "Jump here only if asked how something was done technically.")

# ------------------------------------------------------------------
# A1 — Notebook map
# ------------------------------------------------------------------

s = slide_frame("From survey export to analytical dataset", "Appendix A1",
                "Notebook map", section="ELFI · Appendix")
table(s, M, BODY_TOP, CW,
      [["#", "Notebook", "Question it answers", "Output"],
       [("01", {"mono": True}), "Synthetic data generation",
        "Can I build a publishable twin of the real data?", "synthetic survey dataset"],
       [("02", {"mono": True}), "Exploratory data analysis",
        "Are the data usable for NLP at all?", "cleaned, scored dataset"],
       [("03", {"mono": True}), "Zero-shot text analysis",
        "Is Core Affect feasible without training data?",
        "baseline valence / arousal classes"],
       [("04", {"mono": True}), "German sentiment analysis",
        "Can I get a continuous emotional dimension?",
        ("continuous_valence", {"mono": True})],
       [("05", {"mono": True}), "BERTopic topic modelling",
        "What are employees talking about, and when?",
        "topic id, label, trends, C_v"]],
      col_ratios=[0.06, 0.26, 0.40, 0.28], row_h=0.42, size=11)

tf = card(s, M, 4.9, CW, 1.35)
label(tf, "Why notebooks per question", first=True)
para(tf, "Each notebook exists in a  _high5  variant (real data, confidential) and a public "
         "variant (synthetic data) — same code, different input. That keeps the "
         "confidential analysis reproducible without publishing any comment text.",
     size=11.5, color=INK, line=1.3, space_after=0)
notes(s, "The pipeline is deliberately notebook-based: every notebook answers one "
         "methodological question, and each one can be rerun independently.")

# ------------------------------------------------------------------
# A2 — EDA
# ------------------------------------------------------------------

s = slide_frame("EDA on real High5 data", "Appendix A2",
                "Data understanding", section="ELFI · Appendix")
cols = [("What was analysed",
         ["Rows, columns and schema across 17 survey exports",
          "Availability of free-text comments per iteration",
          "Missing values — real vs. structural vs. “cannot judge”",
          "Distribution across iterations and business units",
          "Duration outliers as a plausibility check"]),
        ("Technical steps",
         ["Harmonise column names across survey generations",
          "Recode Likert answers to numeric 1–5",
          "Reverse-code the workload item so high always means favourable",
          "Separate “cannot judge” from true missing values",
          "Aggregate 54 items into 8 dimension scores",
          "Reliability per dimension (Cronbach α)"]),
        ("Findings",
         ["Several survey generations with differing question blocks — structural "
          "missingness is expected, not an error",
          "α between 0.68 and 0.96; teamwork highest, leadership lowest but acceptable",
          "Commenters rate consistently more critically than non-commenters "
          "(−0.05 to −0.36 points)",
          "≈ 4,300 comments — enough volume for NLP"])]
for i, (head, items) in enumerate(cols):
    tf = card(s, M + i * (col_w + 0.32), BODY_TOP, col_w, 3.2)
    label(tf, head, first=True)
    for j, it in enumerate(items):
        bullet(tf, it, size=10.5, space_after=5 if j < len(items) - 1 else 0)

tf = card(s, M, 5.45, CW, 0.85)
para(tf, "Why it matters: without a validated data base, every downstream model measures "
         "noise. This notebook produces no model — it produces trust in the input.",
     size=12, color=INK, line=1.25, space_after=0, first=True)
notes(s, "Key technical detail to mention if asked: the reverse coding of the workload item, "
         "so that a higher score always means a more favourable assessment across all "
         "dimensions.")

# ------------------------------------------------------------------
# A3 — Real vs synthetic
# ------------------------------------------------------------------

s = slide_frame("How the synthetic dataset was validated", "Appendix A3",
                "Real vs. synthetic", section="ELFI · Appendix")
tf = card(s, M, BODY_TOP, half, 3.15)
label(tf, "Comparison dimensions", first=True)
bullet(tf, "column count, data types, iterations, answer scales",
       bold_prefix="Structure — ", size=11)
bullet(tf, "missingness, comment rate, dimension scores",
       bold_prefix="Distributions — ", size=11)
bullet(tf, "Cronbach α per dimension, inter-dimension correlations",
       bold_prefix="Psychometrics — ", size=11)
bullet(tf, "comment length, vocabulary richness, topic variety",
       bold_prefix="NLP relevance — ", size=11)
bullet(tf, "commenters vs. non-commenters", bold_prefix="Response behaviour — ",
       size=11, space_after=0)

tf = card(s, x2, BODY_TOP, half, 3.15)
label(tf, "Result", first=True)
table(s, x2 + 0.28, BODY_TOP + 0.6, half - 0.56,
      [["Property", "Synthetic"],
       ["Responses × variables", "500 × 66 (v1) → 7,000 (v2)"],
       ["Comment rate", "> 96 %"],
       ["Median comment length", "25 words"],
       ["Dimension scores", "3.1 – 3.7 of 5"],
       ["Cronbach α", "0.68 – 0.96"],
       ["r (teamwork, culture)", "≈ 0.89"],
       ["Business units", "5, n = 82 – 114"]],
      col_ratios=[0.56, 0.44], right_align=(1,), row_h=0.28, size=10)

tf = card(s, M, 5.4, CW, 1.25)
label(tf, "Honest limitation", first=True, color=CRIT)
para(tf, "Structure and statistics transfer well; language does not. Synthetic comments are "
         "more repetitive and share phrasing, which makes topic modelling measurably harder "
         "than on real text. The generator was reworked after the first EDA for exactly that "
         "reason — and the topic metrics still carry the effect.",
     size=11.5, color=INK, line=1.3, space_after=0)
notes(s, "This is the slide to use if someone challenges the topic-model numbers: part of "
         "the weakness is a property of the synthetic corpus, not of the method.")

# ------------------------------------------------------------------
# A4 — Zero-shot
# ------------------------------------------------------------------

s = slide_frame("Zero-shot prototype — and why it was replaced", "Appendix A4",
                "Baseline model", section="ELFI · Appendix")
cols = [("Setup",
         ["Model  MoritzLaurer/mDeBERTa-v3-base-mnli-xnli",
          "Classification via natural language inference — no training data needed",
          "Valence labels: positive / neutral / negative",
          "Arousal labels: high / medium / low activation",
          "Both dimensions scored per comment, exported as Core Affect coordinates"]),
        ("What it proved",
         ["End-to-end pipeline works: load → classify → derive coordinates "
          "→ export",
          "The Core Affect idea is technically feasible on German comments",
          "The only component so far that produces arousal estimates at all",
          "Useful as a fast feasibility test before investing in a specialised model"]),
        ("Why it stays a baseline",
         ["Unstable arousal — label wording drives the result, and three coarse classes "
          "are not enough",
          "Discrete output — no continuous dimension for the Core Affect Map",
          "Slow — ≈ 1.7 comments/s, about 50 min for 5,000 comments, roughly 12× "
          "slower than the specialised sentiment model"])]
for i, (head, items) in enumerate(cols):
    lift = i == 2
    tf = card(s, M + i * (col_w + 0.32), BODY_TOP, col_w, 3.15, lift=lift)
    label(tf, head, first=True)
    for j, it in enumerate(items):
        bullet(tf, it, size=10.5, space_after=5 if j < len(items) - 1 else 0)

pill(s, M + 2 * (col_w + 0.32), 5.4, 2.6, "kept as baseline, not product", CRIT)
tf = textbox(s, M, 5.4, half + 1.0, 0.9)
para(tf, "Decision: the zero-shot model remains documented as a proof of concept. The "
         "product uses German Sentiment BERT for valence, and arousal is being solved "
         "separately (see A9).",
     size=11.5, color=MUTED, line=1.25, space_after=0, first=True)
notes(s, "If asked why not just keep zero-shot: it is 12 times slower, produces discrete "
         "classes instead of a continuous scale, and its arousal output shifts when the "
         "candidate labels are reworded — not defensible as a measurement.")

# ------------------------------------------------------------------
# A5 — Sentiment
# ------------------------------------------------------------------

s = slide_frame("German Sentiment BERT → continuous valence", "Appendix A5",
                "Sentiment layer", section="ELFI · Appendix")
tf = card(s, M, BODY_TOP, half, 3.25)
label(tf, "Technical steps", first=True)
bullet(tf, "Filter to valid comments (has_comment == True) — 5,000 of 7,000 rows",
       size=11)
bullet(tf, "Model  oliverguhr/german-sentiment-bert  via HF pipeline, top_k=None, "
           "truncation enabled", size=11)
bullet(tf, "Keep all three probabilities per comment, not just the argmax label", size=11)
bullet(tf, "Derive  continuous_valence = P(pos) − P(neg),  range −1 … +1", size=11)
bullet(tf, "Merge back onto the survey rows, export the enriched dataset for Power BI",
       size=11)
bullet(tf, "Throughput ≈ 20 comments/s — 5,000 comments in about 4 minutes",
       size=11, space_after=0)

tf = card(s, x2, BODY_TOP, half, 3.25)
label(tf, "From version 1 to version 2", first=True)
table(s, x2 + 0.28, BODY_TOP + 0.6, half - 0.56,
      [["Zero-shot", "German Sentiment BERT"],
       ["3 discrete states", "continuous −1 … +1"],
       ["low differentiation", "intensity is preserved"],
       ["limited dashboard use", "axis 1 of the Core Affect Map"]],
      col_ratios=[0.45, 0.55], row_h=0.3, size=10.5)
tf2 = textbox(s, x2 + 0.28, BODY_TOP + 2.1, half - 0.56, 1.2)
para(tf2, "The valence distribution is polarised, not centred: mass at both ends, mean "
          "≈ −0.24 in the exported sample. That matches the EDA finding that people "
          "who comment are the more critical subgroup.",
     size=11.5, color=INK, line=1.28, space_after=0, first=True)

tf = card(s, M, 5.45, CW, 1.25)
label(tf, "Open point", first=True, color=WARM)
para(tf, "The model is trained on general German, not on survey language. Domain fit still "
         "needs a manual spot check against a hand-labelled sample of comments before the "
         "valence axis is treated as validated.",
     size=11.5, color=INK, line=1.28, space_after=0)
notes(s, "The key modelling decision here: keeping the full probability vector instead of the "
         "predicted label. That is what makes a continuous axis possible, and it is what the "
         "Core Affect Map needs.")

# ------------------------------------------------------------------
# A6 — BERTopic representations
# ------------------------------------------------------------------

s = slide_frame("BERTopic: four representation experiments", "Appendix A6",
                "Topic layer — iterations", section="ELFI · Appendix")
tf = card(s, M, BODY_TOP, half, 2.5)
label(tf, "Pipeline", first=True)
bullet(tf, "paraphrase-multilingual-MiniLM-L12-v2", bold_prefix="Embeddings —  ",
       size=11)
bullet(tf, "dimensionality reduction, cosine metric, 5 components",
       bold_prefix="UMAP —  ", size=11)
bullet(tf, "density clustering, EOM cluster selection", bold_prefix="HDBSCAN —  ",
       size=11)
bullet(tf, "topic keyword extraction", bold_prefix="c-TF-IDF + representation model —  ",
       size=11)
bullet(tf, "fixed random_state = 42 throughout, so configurations stay comparable",
       size=11, space_after=0)

tf = card(s, x2, BODY_TOP, half, 2.5)
label(tf, "What was changed, and what happened", first=True)
table(s, x2 + 0.28, BODY_TOP + 0.6, half - 0.56,
      [["Configuration", "Topics"],
       ["Baseline BERTopic", "27"],
       ["+ CountVectorizer, extended stopwords", "42"],
       ["+ KeyBERTInspired representation", "38"],
       [("+ BM25 weighting — rejected", {"color": CRIT}),
        ("collapsed", {"color": CRIT})]],
      col_ratios=[0.72, 0.28], right_align=(1,), row_h=0.3, size=10.5)

tf = card(s, M, 4.85, CW, 1.9)
label(tf, "Reading the experiments", first=True)
bullet(tf, "Survey boilerplate (“diskutiert”, “Kontext”, "
           "“wird”) dominated the keyword lists. An extended German + "
           "survey-specific stopword list, refined over several passes, mattered most for "
           "readability.", bold_prefix="Stopwords —  ", size=11)
bullet(tf, "re-ranks candidate keywords by embedding similarity to the topic instead of "
           "raw frequency — noticeably more interpretable labels.",
       bold_prefix="KeyBERTInspired —  ", size=11)
bullet(tf, "expected to sharpen topic-specific terms; instead nearly all comments collapsed "
           "into one dominant topic with a few small satellites. Dropped, and documented as "
           "a negative result.", bold_prefix="BM25 —  ", size=11, space_after=0)
notes(s, "Good place to show method: each change was evaluated on its own before the next was "
         "added, and the BM25 result is reported even though it failed.")

# ------------------------------------------------------------------
# A7 — Grid searches
# ------------------------------------------------------------------

s = slide_frame("Two grid searches: clustering, then projection", "Appendix A7",
                "Topic layer — hyperparameters", section="ELFI · Appendix")
tf = card(s, M, BODY_TOP, half, 4.0)
label(tf, "Step 1 — HDBSCAN min_cluster_size", first=True)
para(tf, "5 configurations · synthetic data, n = 5,000",
     size=9.5, font=MONO, color=FAINT, caps=True, space_after=6)
table(s, M + 0.28, BODY_TOP + 0.95, half - 0.56,
      [["min_cluster_size", "Topics", "Outliers", "Rate"],
       [("5", {"bold": True, "color": ACCENT}), ("208", {"bold": True, "color": ACCENT}),
        ("1,045", {"bold": True, "color": ACCENT}), ("20.9 %", {"bold": True,
                                                                "color": ACCENT})],
       ["10", "101", "878", "17.6 %"],
       ["15", "68", "863", "17.3 %"],
       ["20", "37", "556", "11.1 %"],
       ["25", "22", "108", "2.2 %"]],
      col_ratios=[0.34, 0.2, 0.24, 0.22], right_align=(1, 2, 3), row_h=0.28, size=10.5)
tf2 = textbox(s, M + 0.28, BODY_TOP + 2.85, half - 0.56, 1.3)
para(tf2, "Clear trade-off: larger minimum clusters crush the outlier rate but also crush "
          "granularity — 22 topics for 5,000 comments says almost nothing. "
          "min_cluster_size = 5 was carried forward to keep a differentiated structure, with "
          "the outlier rate to be addressed by the projection step instead.",
     size=11, color=INK, line=1.25, space_after=0, first=True)

tf = card(s, x2, BODY_TOP, half, 4.0)
label(tf, "Step 2 — UMAP n_neighbors × min_dist", first=True)
para(tf, "9 combinations · min_cluster_size fixed at 5",
     size=9.5, font=MONO, color=FAINT, caps=True, space_after=6)
table(s, x2 + 0.28, BODY_TOP + 0.95, half - 0.56,
      [["n_neighbors", "min_dist", "Topics", "Outlier rate"],
       ["15", "0.3", "14", "0.4 %"],
       ["30", "0.3", "16", "0.8 %"],
       ["5", "0.3", "95", "10.6 %"],
       [("5", {"bold": True, "color": ACCENT}), ("0.0", {"bold": True, "color": ACCENT}),
        ("304", {"bold": True, "color": ACCENT}), ("18.1 %", {"bold": True,
                                                              "color": ACCENT})],
       ["15", "0.0", "208", "20.9 %"],
       ["30", "0.1", "168", "27.5 %"]],
      col_ratios=[0.28, 0.22, 0.2, 0.3], right_align=(1, 2, 3), row_h=0.28, size=10.5)
tf2 = textbox(s, x2 + 0.28, BODY_TOP + 3.15, half - 0.56, 1.1)
para(tf2, "6 of 9 shown; the remaining combinations landed between 22 % and 25 %. A high "
          "min_dist looks excellent on the outlier metric and is misleading — 14 topics "
          "for 5,000 comments is one large blob, not a finding.",
     size=11, color=INK, line=1.25, space_after=0, first=True)
notes(s, "The point to make: the metric alone would have selected the wrong model. "
         "min_dist = 0.3 gives a 0.4 % outlier rate and is useless. Metrics were read "
         "together with topic counts and manual inspection of the keyword lists.")

# ------------------------------------------------------------------
# A8 — Final model & coherence
# ------------------------------------------------------------------

s = slide_frame("Final configuration and its honest verdict", "Appendix A8",
                "Topic layer — evaluation", section="ELFI · Appendix")
tf = card(s, M, BODY_TOP, col_w, 3.6, lift=True)
label(tf, "Selected configuration", first=True)
bullet(tf, "CountVectorizer with extended German + survey stopwords", size=10.5)
bullet(tf, "KeyBERTInspired representation", size=10.5)
bullet(tf, "UMAP  n_neighbors = 5,  min_dist = 0.0,  5 components, cosine", size=10.5)
bullet(tf, "HDBSCAN  min_cluster_size = 5,  EOM selection", size=10.5)
bullet(tf, "calculate_probabilities = True, so each comment carries a topic probability",
       size=10.5, space_after=0)

tf = card(s, M + col_w + 0.32, BODY_TOP, col_w, 3.6)
label(tf, "Quality — synthetic run", first=True)
table(s, M + col_w + 0.6, BODY_TOP + 0.6, col_w - 0.56,
      [["Documents", "5,000"],
       ["Topics excl. outlier", "303"],
       ["Outliers", "905  (18.1 %)"],
       ["Mean / median size", "13.5 / 11"],
       ["Topics ≥ 25 comments", "22"],
       ["Topics < 10 comments", "115"],
       [("C_v coherence", {"bold": True}), ("0.350", {"bold": True, "color": ACCENT})]],
      header=False, col_ratios=[0.56, 0.44], right_align=(1,), row_h=0.28, size=10.5)
tf2 = textbox(s, M + col_w + 0.6, BODY_TOP + 2.7, col_w - 0.56, 0.7)
para(tf2, "Real-data run for comparison: 4,297 documents → 215 topics, 29.5 % outliers.",
     size=10.5, color=MUTED, line=1.25, space_after=0, first=True)

tf = card(s, M + 2 * (col_w + 0.32), BODY_TOP, col_w, 3.6)
label(tf, "Verdict & fix", first=True, color=CRIT)
bullet(tf, "303 topics with a median of 11 comments are micro-topics; only 22 are large "
           "enough to act on.", bold_prefix="Too granular.  ", size=10.5)
bullet(tf, "Partly the synthetic text: repetitive phrasing produces overlapping keyword "
           "sets across topics.", bold_prefix="C_v = 0.350 is modest.  ", size=10.5)
bullet(tf, "reduce topics to a manageable set, re-evaluate coherence per configuration, "
           "then repeat the whole comparison on real comments.",
       bold_prefix="Next:  ", size=10.5, space_after=0)

tf = textbox(s, M, 5.9, CW, 0.7)
para(tf, "Reported as measured. The coherence value is not yet where it needs to be, and the "
         "granularity problem is named as a risk rather than smoothed over.",
     size=11.5, color=MUTED, italic=True, line=1.25, space_after=0, first=True)
notes(s, "Numbers come straight from the notebook quality check. If asked why the real-data "
         "run has a higher outlier rate: fewer documents, far more heterogeneous language.")

# ------------------------------------------------------------------
# A9 — Arousal
# ------------------------------------------------------------------

s = slide_frame("How arousal gets its axis", "Appendix A9",
                "Open question", section="ELFI · Appendix")
tf = card(s, M, BODY_TOP, col_w, 2.95)
label(tf, "The gap", first=True)
para(tf, "German Sentiment BERT delivers valence only. Without arousal the Core Affect Map "
         "has one axis — stress and resignation stay indistinguishable, although they "
         "call for opposite measures.",
     size=11.5, color=INK, line=1.3, space_after=6)
para(tf, "The zero-shot prototype produced arousal classes (low / medium / high), but they "
         "were not stable enough to build a dashboard axis on.",
     size=11.5, color=MUTED, line=1.3, space_after=0)

tf = card(s, M + col_w + 0.32, BODY_TOP, col_w, 2.95, lift=True)
label(tf, "Preferred route — emotion mapping", first=True)
bullet(tf, "A German emotion classifier (e.g. ChrisLalk/German-Emotions) predicts discrete "
           "emotions: frustration, uncertainty, pride, hope, anger, joy", size=10.5)
bullet(tf, "Each emotion carries known valence and arousal coordinates from affect research",
       size=10.5)
bullet(tf, "Probability-weighted mapping yields continuous coordinates per comment",
       size=10.5)
bullet(tf, "Reuses established theory instead of inventing a scale", size=10.5,
       space_after=0)

tf = card(s, M + 2 * (col_w + 0.32), BODY_TOP, col_w, 2.95)
label(tf, "Fallback & validation", first=True)
bullet(tf, "Train a custom arousal model on emotional language, comment features and survey "
           "items — needs labelled data, only worth it with time to spare", size=10.5)
bullet(tf, "Either way: correlate valence and arousal against engagement, team, culture and "
           "leadership scores", size=10.5)
bullet(tf, "If the emotional layer carries signal, those relationships have to show up",
       size=10.5, space_after=0)

tf = card(s, M, 5.25, CW, 1.3)
label(tf, "Why this matters for the product", first=True)
para(tf, "Valence alone answers “is it good or bad?”. Arousal answers “is "
         "there energy behind it?” — the difference between a team that is angry "
         "and a team that has given up. The recommended measures for those two states are "
         "opposites, which is the entire argument for the Barrett-inspired framework.",
     size=11.5, color=INK, line=1.3, space_after=0)
notes(s, "This is the most likely question from the examiners: how do you get arousal? Answer "
         "with the mapping route, and be explicit that it is not implemented yet.")

# ------------------------------------------------------------------
# A10 — Performance
# ------------------------------------------------------------------

s = slide_frame("Performance and reproducibility", "Appendix A10",
                "Engineering", section="ELFI · Appendix")
tf = card(s, M, BODY_TOP, half, 3.2)
label(tf, "Problem", first=True)
para(tf, "Several steps are expensive and are recomputed on every notebook run:",
     size=11.5, color=INK, line=1.25, space_after=6)
bullet(tf, "Sentence embeddings for 5,000 comments — ≈ 30 s per fit, repeated for "
           "every grid cell", size=11)
bullet(tf, "Sentiment inference — ≈ 4 min per full pass", size=11)
bullet(tf, "Zero-shot classification — ≈ 50 min per full pass", size=11)
bullet(tf, "BERTopic — 5 HDBSCAN + 9 UMAP configurations during tuning", size=11,
       space_after=0)

tf = card(s, x2, BODY_TOP, half, 3.2)
label(tf, "Options evaluated", first=True)
table(s, x2 + 0.28, BODY_TOP + 0.6, half - 0.56,
      [["Format", "Assessment"],
       ["CSV", "Readable, but large, slow, loses dtypes"],
       ["SQLite", "Good for relational queries across iterations"],
       [("Parquet", {"bold": True, "color": ACCENT}),
        ("Columnar, compressed, keeps dtypes, fastest reload", {"color": ACCENT})]],
      col_ratios=[0.24, 0.76], row_h=0.34, size=10.5)
tf2 = textbox(s, x2 + 0.28, BODY_TOP + 2.05, half - 0.56, 1.0)
para(tf2, "Current plan: persist cleaned_data.parquet, sentiment_results.parquet and "
          "topic_results.parquet, keyed by comment id so the layers can be joined "
          "independently.", size=11, color=INK, line=1.25, space_after=0, first=True)

pill(s, M, 5.45, 2.4, "evaluated, not implemented", WARM)
tf = textbox(s, M, 5.85, CW, 0.7)
para(tf, "Say this plainly if it comes up: caching is a decision on the table, not a "
         "finished part of the pipeline. The reproducibility argument — same input, "
         "same intermediate artefacts, same dashboard — is why it is worth doing before "
         "the final evaluation runs.",
     size=11.5, color=MUTED, line=1.25, space_after=0, first=True)
notes(s, "Do not oversell this. It is currently under evaluation, and Parquet is the "
         "pragmatic choice for an analytical pandas workflow.")

# ------------------------------------------------------------------
# A11 — Definition of done
# ------------------------------------------------------------------

s = slide_frame("What makes this project a success", "Appendix A11",
                "Definition of done", section="ELFI · Appendix")
blocks = [("Data product",
           ["≥ 95 % of comments processed automatically",
            "Reproducible, rerunnable pipeline",
            "Dashboard-ready export generated automatically",
            "Power BI integration"]),
          ("NLP",
           ["Continuous valence — done",
            "Interpretable topic clusters — granularity open",
            "Continuous arousal — open"]),
          ("Predictive layer",
           ["Comment → engagement score",
            "TF-IDF baseline vs. gradient boosting",
            "Targets: MAE < 0.60, RMSE < 0.80, R² > 0.30"]),
          ("Business value",
           ["3–5 actionable findings",
            "e.g. valence ↔ engagement, workload topics ↔ low sentiment",
            "Measurable differences between business units"])]
for i, (head, items) in enumerate(blocks):
    tf = card(s, M + i * (kw + 0.28), BODY_TOP, kw, 2.5)
    label(tf, head, first=True)
    for j, it in enumerate(items):
        bullet(tf, it, size=10.5, space_after=5 if j < len(items) - 1 else 0)

rect(s, M, 4.95, 0.03, 1.15, ACCENT)
tf = textbox(s, M + 0.28, 5.0, CW - 0.5, 1.1)
para(tf, "The project succeeds if it shows a working end-to-end NLP pipeline, a "
         "Barrett-inspired Core Affect dashboard, meaningful topics, and insights a manager "
         "can actually act on.",
     size=15, font=SERIF, color=INK, line=1.25, space_after=0, first=True)
notes(s, "Success is deliberately not defined by model performance alone — the project "
         "combines data engineering, NLP, psychology and dashboard design.")

# ------------------------------------------------------------------
# Page numbers + save
# ------------------------------------------------------------------

total = len(prs.slides._sldIdLst)
for n, slide in enumerate(prs.slides, start=1):
    if n == 1:
        continue
    stamp_page(slide, n, total)

prs.core_properties.title = "ELFI — Emotional & Latent Feedback Interpretation"
prs.core_properties.subject = "Capstone midterm presentation"
prs.core_properties.comments = ("Midterm deck: main story (slides 1-9) plus technical "
                                "appendix (A1-A11). Speaker notes in the notes pane.")

OUT = "/Users/idetemple/Desktop/fish_happens/ELFI-capstone/presentation/02_Midterm_Presentation.pptx"
prs.save(OUT)
print(f"saved {OUT}  — {total} slides")
